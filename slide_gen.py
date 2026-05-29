"""
uruvagam.slide_gen
~~~~~~~~~~~~~~~~~~~
Generates a themed PPTX from structured content JSON.
16:9 widescreen format, configurable branding, speaker notes included.
"""

import argparse
import json
import os
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.enum.text import PP_ALIGN

from theme_config import DEFAULT_THEME_NAME, hex_to_rgb, load_theme

# ── Theme colours (dark corporate) ────────────────────────────────────────────
C = {
    "primary":    RGBColor(0x1A, 0x1A, 0x2E),  # deep navy bg
    "surface":    RGBColor(0x16, 0x21, 0x3E),  # card/header bg
    "elevated":   RGBColor(0x0F, 0x3D, 0x6E),  # raised element
    "cyan":       RGBColor(0x53, 0xC4, 0xFF),  # accent / highlight
    "cyan_dim":   RGBColor(0x2A, 0x7A, 0xAA),  # subdued accent
    "white":      RGBColor(0xFF, 0xFF, 0xFF),
    "off_white":  RGBColor(0xD8, 0xDF, 0xEB),  # body text
    "muted":      RGBColor(0x8A, 0x9A, 0xBB),  # secondary text
    "number_bg":  RGBColor(0x0A, 0x2A, 0x50),  # number badge fill
}

W = Inches(16)
H = Inches(9)


def generate_pptx(
    content: dict,
    output_path: str,
    theme: str | None = DEFAULT_THEME_NAME,
    template_path: str | None = None,
    logo_path: str | None = None,
) -> str:
    """Main entry point. Returns path to saved .pptx file."""
    theme_cfg = load_theme(theme)
    colors = _ppt_colors(theme_cfg)
    logo = _optional_path(logo_path)

    if template_path:
        template = Path(template_path)
        if not template.exists():
            raise FileNotFoundError(f"template not found: {template}")
        prs = Presentation(str(template))
        _remove_all_slides(prs)
    else:
        prs = Presentation()
        prs.slide_width = W
        prs.slide_height = H

    slides = content.get("slides", [])

    # Always generate title slide from content metadata
    if template_path:
        _add_template_title_slide(prs, content, theme_cfg, colors, logo)
    else:
        _add_title_slide(prs, content, theme_cfg, colors, logo)

    # Generate remaining slides
    for slide_data in slides:
        stype = slide_data.get("slide_type", "content")
        if stype == "title":
            continue  # already added above
        if template_path:
            if stype in ("objectives",):
                _add_template_content_slide(prs, slide_data, theme_cfg, colors, logo, content.get("objectives", []))
            elif stype == "summary":
                _add_template_content_slide(prs, slide_data, theme_cfg, colors, logo)
            else:
                _add_template_content_slide(prs, slide_data, theme_cfg, colors, logo)
        elif stype in ("objectives",):
            _add_objectives_slide(prs, slide_data, content.get("objectives", []), theme_cfg, colors, logo)
        elif stype == "summary":
            _add_summary_slide(prs, slide_data, theme_cfg, colors, logo)
        else:
            _add_content_slide(prs, slide_data, theme_cfg, colors, logo)

    Path(output_path).parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)
    return output_path


# ── Template slide builders ────────────────────────────────────────────────────

def _add_template_title_slide(prs: Presentation, content: dict, theme: dict, colors: dict, logo_path: Path | None):
    layout = _layout(prs, theme["ppt"].get("layout_title", "Title Slide White"), fallback=0)
    slide = prs.slides.add_slide(layout)
    _fill_bg(slide, colors["primary"])

    title = _find_placeholder(slide, PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE)
    subtitle = _find_placeholder(slide, PP_PLACEHOLDER.SUBTITLE)
    duration = content.get("duration_minutes", 10)
    audience = content.get("_audience", "Engineering Team")
    # _tagline/_presenter override the default "Training Module - ..." subtitle when present
    if content.get("_tagline") or content.get("_presenter"):
        subtitle_text = "\n".join(
            t for t in (content.get("_tagline", ""), content.get("_presenter", "")) if t
        )
    else:
        subtitle_text = f"Training Module - {duration} minutes - {audience}"

    if title:
        _set_shape_text(
            title,
            content.get("title", "Training"),
            size=Pt(42),
            bold=True,
            color=colors["white"],
            font_name=theme["fonts"].get("heading"),
        )
    else:
        _text_box(
            slide, Inches(0.75), Inches(2.2), Inches(13.8), Inches(1.6),
            content.get("title", "Training"), size=Pt(42), bold=True,
            color=colors["white"], wrap=True, font_name=theme["fonts"].get("heading"),
        )

    if subtitle:
        _set_shape_text(
            subtitle,
            subtitle_text,
            size=Pt(18),
            color=colors["muted"],
            font_name=theme["fonts"].get("body"),
        )
    else:
        _text_box(
            slide, Inches(0.75), Inches(4.5), Inches(12), Inches(0.9),
            subtitle_text,
            size=Pt(18), color=colors["muted"], font_name=theme["fonts"].get("body"),
        )

    _add_footer(slide, theme, colors, logo_path)
    _set_notes(slide, _title_notes(content))


def _add_template_content_slide(
    prs: Presentation,
    slide_data: dict,
    theme: dict,
    colors: dict,
    logo_path: Path | None,
    override_bullets: list | None = None,
):
    stype = slide_data.get("slide_type", "content")
    layout_name = theme["ppt"].get("layout_summary" if stype == "summary" else "layout_content", "Title and Content")
    layout = _layout(prs, layout_name, fallback=1)
    slide = prs.slides.add_slide(layout)
    _fill_bg(slide, colors["primary"])

    title = _find_placeholder(slide, PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE)
    body = _find_placeholder(slide, PP_PLACEHOLDER.OBJECT, PP_PLACEHOLDER.BODY)
    title_text = "Learning Objectives" if stype == "objectives" else slide_data.get("title", "")
    bullets = override_bullets or slide_data.get("bullets", [])

    if title:
        _set_shape_text(
            title,
            title_text,
            size=Pt(28),
            bold=True,
            color=colors["cyan"],
            font_name=theme["fonts"].get("heading"),
        )
    else:
        _header_bar(slide, title_text, colors, theme)

    diagram = slide_data.get("diagram", "")
    # when diagram present, compress bullet rows to leave room
    max_bullets = 5 if diagram else 8
    row_height = 0.78 if (len(bullets) > 6 or diagram) else 1.02

    if body:
        _set_bullets(body, bullets[:max_bullets], colors, theme, subtitle=slide_data.get("subtitle", ""))
        y_start = Inches(1.65)
        subtitle = slide_data.get("subtitle", "")
        if subtitle:
            _text_box(
                slide, Inches(0.65), Inches(1.3), Inches(14), Inches(0.45),
                subtitle, size=Pt(15), color=colors["muted"], font_name=theme["fonts"].get("body"),
            )
            y_start = Inches(1.95)
        for i, bullet in enumerate(bullets[:max_bullets]):
            _bullet_row(slide, i, bullet, y_start + Inches(i * row_height), colors, theme)

    # diagram block below bullets (PPTX: fixed-width text box with monospace font)
    if diagram:
        diag_y = Inches(1.65) + Inches(max_bullets * row_height) + Inches(0.1)
        diag_h = Inches(7.5 - 1.65 - max_bullets * row_height - 0.2)
        if diag_h > Inches(0.8):
            tb = slide.shapes.add_textbox(Inches(0.5), diag_y, Inches(14.5), diag_h)
            tf = tb.text_frame
            tf.word_wrap = False
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = diagram.strip("\n")
            run.font.size = Pt(7)
            run.font.name = "Courier New"
            run.font.color.rgb = colors["cyan_dim"]

    _add_footer(slide, theme, colors, logo_path)
    _set_notes(slide, slide_data.get("speaker_notes", ""))


# ── Slide builders ─────────────────────────────────────────────────────────────

def _add_title_slide(prs: Presentation, content: dict, theme: dict, colors: dict, logo_path: Path | None):
    slide = _blank_slide(prs)
    _fill_bg(slide, colors["primary"])

    # Left cyan bar
    _rect(slide, 0, 0, Inches(0.18), H, colors["cyan"])

    # Large title
    _text_box(
        slide, Inches(0.45), Inches(2.2), Inches(14.5), Inches(2.4),
        content.get("title", "Training"), size=Pt(46), bold=True,
        color=colors["white"], wrap=True, font_name=theme["fonts"].get("heading"),
    )

    # Subtitle / module info
    duration = content.get("duration_minutes", 10)
    audience = content.get("_audience", "Engineering Team")
    _text_box(
        slide, Inches(0.45), Inches(5.0), Inches(10), Inches(0.7),
        f"Training Module - {duration} minutes",
        size=Pt(20), color=colors["cyan"], font_name=theme["fonts"].get("body"),
    )

    # Bottom rule
    _rect(slide, Inches(0.45), Inches(6.2), Inches(13), Inches(0.04), colors["cyan_dim"])

    # Tagline
    _text_box(
        slide, Inches(0.45), Inches(6.5), Inches(10), Inches(0.6),
        theme["ppt"].get("footer_text", "Generated with uruvagam"),
        size=Pt(14), color=colors["muted"], font_name=theme["fonts"].get("body"),
    )

    _add_footer(slide, theme, colors, logo_path)
    _set_notes(slide, _title_notes(content))


def _add_objectives_slide(prs: Presentation, slide_data: dict, objectives: list, theme: dict, colors: dict, logo_path: Path | None):
    slide = _blank_slide(prs)
    _fill_bg(slide, colors["primary"])
    _header_bar(slide, "Learning Objectives", colors, theme)

    items = objectives or slide_data.get("bullets", [])
    for i, obj in enumerate(items[:6]):
        y = Inches(1.6 + i * 1.12)
        _bullet_row(slide, i, obj, y, colors, theme)

    _add_footer(slide, theme, colors, logo_path)
    _set_notes(slide, slide_data.get("speaker_notes", ""))


def _add_content_slide(prs: Presentation, slide_data: dict, theme: dict, colors: dict, logo_path: Path | None):
    slide = _blank_slide(prs)
    _fill_bg(slide, colors["primary"])
    _header_bar(slide, slide_data.get("title", ""), colors, theme)

    subtitle = slide_data.get("subtitle", "")
    if subtitle:
        _text_box(slide, Inches(0.4), Inches(1.4), Inches(15), Inches(0.5),
                  subtitle, size=Pt(16), color=colors["muted"], font_name=theme["fonts"].get("body"))

    y_start = Inches(1.65) if not subtitle else Inches(2.0)
    bullets = slide_data.get("bullets", [])
    for i, bullet in enumerate(bullets[:6]):
        y = y_start + Inches(i * 1.1)
        _bullet_row(slide, i, bullet, y, colors, theme)

    _add_footer(slide, theme, colors, logo_path)
    _set_notes(slide, slide_data.get("speaker_notes", ""))


def _add_summary_slide(prs: Presentation, slide_data: dict, theme: dict, colors: dict, logo_path: Path | None):
    slide = _blank_slide(prs)
    _fill_bg(slide, colors["primary"])

    # Gradient-like top accent
    _rect(slide, 0, 0, W, Inches(0.18), colors["cyan"])
    _rect(slide, 0, Inches(0.18), W, Inches(1.2), colors["elevated"])

    _text_box(
        slide, Inches(0.4), Inches(0.25), Inches(15), Inches(1.0),
        slide_data.get("title", "Summary & Next Steps"),
        size=Pt(30), bold=True, color=colors["cyan"], font_name=theme["fonts"].get("heading"),
    )

    bullets = slide_data.get("bullets", [])
    for i, b in enumerate(bullets[:5]):
        y = Inches(1.7 + i * 1.1)
        _bullet_row(slide, i, b, y, colors, theme, accent_color=colors["cyan"])

    # CTA box at bottom
    _rect(slide, Inches(0.4), Inches(7.8), Inches(15.2), Inches(0.9), colors["elevated"])
    _text_box(
        slide, Inches(0.6), Inches(7.88), Inches(14), Inches(0.6),
        "Questions? Connect with the team - slides available in your knowledge base",
        size=Pt(14), color=colors["muted"], font_name=theme["fonts"].get("body"),
    )

    _add_footer(slide, theme, colors, logo_path)
    _set_notes(slide, slide_data.get("speaker_notes", ""))


# ── Reusable primitives ────────────────────────────────────────────────────────

def _blank_slide(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])  # blank


def _fill_bg(slide, color: RGBColor):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _rect(slide, x, y, w, h, color: RGBColor):
    shape = slide.shapes.add_shape(1, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def _text_box(slide, x, y, w, h, text, size=Pt(20), bold=False,
              color=None, wrap=True, align=PP_ALIGN.LEFT, font_name=None):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    run = p.runs[0]
    run.font.size = size
    run.font.bold = bold
    if font_name:
        run.font.name = font_name
    if color:
        run.font.color.rgb = color
    return tb


def _header_bar(slide, title_text: str, colors: dict, theme: dict):
    """Dark header band with title text and left cyan bar."""
    _rect(slide, 0, 0, Inches(0.18), H, colors["cyan"])
    _rect(slide, Inches(0.18), 0, W - Inches(0.18), Inches(1.3), colors["surface"])
    _text_box(
        slide, Inches(0.4), Inches(0.12), Inches(15), Inches(1.1),
        title_text, size=Pt(30), bold=True, color=colors["cyan"],
        font_name=theme["fonts"].get("heading"),
    )
    # Thin rule below header
    _rect(slide, Inches(0.4), Inches(1.3), Inches(15.2), Inches(0.025), colors["cyan_dim"])


def _bullet_row(slide, index: int, text: str, y, colors: dict, theme: dict, accent_color=None):
    """Numbered bullet row with background pill."""
    accent = accent_color or colors["cyan"]
    # Row background (subtle)
    row_color = colors["number_bg"] if index % 2 == 0 else colors["row_even"]
    _rect(slide, Inches(0.28), y, Inches(15.3), Inches(0.88), row_color)

    # Number badge
    badge = slide.shapes.add_textbox(Inches(0.35), y + Inches(0.12), Inches(0.5), Inches(0.65))
    btf = badge.text_frame
    bp = btf.paragraphs[0]
    bp.text = str(index + 1)
    bp.alignment = PP_ALIGN.CENTER
    br = bp.runs[0]
    br.font.size = Pt(15)
    br.font.bold = True
    br.font.color.rgb = accent
    br.font.name = theme["fonts"].get("body")

    # Bullet text
    _text_box(
        slide, Inches(0.95), y + Inches(0.08), Inches(14.4), Inches(0.75),
        text, size=Pt(20), color=colors["off_white"], wrap=True,
        font_name=theme["fonts"].get("body"),
    )


def _set_notes(slide, text: str):
    if text:
        slide.notes_slide.notes_text_frame.text = text


def _ppt_colors(theme: dict) -> dict:
    colors = theme["colors"]
    return {
        "primary": _rgb(colors["slide_background"]),
        "surface": _rgb(colors["surface"]),
        "elevated": _rgb(colors["elevated"]),
        "cyan": _rgb(colors["accent"]),
        "cyan_dim": _rgb(colors["accent_2"]),
        "white": _rgb(colors["text"]),
        "off_white": _rgb(colors["body_text"]),
        "muted": _rgb(colors["muted"]),
        "number_bg": _rgb(colors["row_odd"]),
        "row_even": _rgb(colors["row_even"]),
    }


def _rgb(value: str) -> RGBColor:
    return RGBColor(*hex_to_rgb(value))


def _optional_path(value: str | None) -> Path | None:
    if not value:
        return None
    path = Path(value)
    if not path.exists():
        raise FileNotFoundError(f"logo not found: {path}")
    return path


def _title_notes(content: dict) -> str:
    for s in content.get("slides", []):
        if s.get("slide_type") == "title":
            return s.get("speaker_notes", "")
    return f"Welcome to this training on {content.get('title', '')}."


def _remove_all_slides(prs: Presentation):
    slide_ids = prs.slides._sldIdLst
    for slide_id in list(slide_ids):
        prs.part.drop_rel(slide_id.rId)
        slide_ids.remove(slide_id)


def _layout(prs: Presentation, name: str, fallback: int):
    for layout in prs.slide_layouts:
        if layout.name == name:
            return layout
    return prs.slide_layouts[min(fallback, len(prs.slide_layouts) - 1)]


def _find_placeholder(slide, *types):
    for shape in slide.placeholders:
        if shape.placeholder_format.type in types:
            return shape
    return None


def _set_shape_text(shape, text: str, size=Pt(20), bold=False, color=None, font_name=None):
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = str(text)
    run.font.size = size
    run.font.bold = bold
    if font_name:
        run.font.name = font_name
    if color:
        run.font.color.rgb = color


def _set_bullets(shape, bullets: list, colors: dict, theme: dict, subtitle: str = ""):
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    if subtitle:
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.name = theme["fonts"].get("body")
        p.font.color.rgb = colors["muted"]
        start_index = 1
    else:
        start_index = 0

    for i, bullet in enumerate(bullets[:6]):
        p = tf.paragraphs[0] if i == 0 and start_index == 0 else tf.add_paragraph()
        p.text = str(bullet)
        p.level = 0
        p.font.size = Pt(21)
        p.font.name = theme["fonts"].get("body")
        p.font.color.rgb = colors["off_white"]
        p.space_after = Pt(10)


def _add_footer(slide, theme: dict, colors: dict, logo_path: Path | None):
    show_footer = logo_path is not None or theme.get("name") != DEFAULT_THEME_NAME
    if not show_footer:
        return

    x = Inches(0.35)
    y = Inches(8.15)
    if logo_path:
        slide.shapes.add_picture(
            str(logo_path),
            x,
            y,
            width=Inches(float(theme["ppt"].get("logo_width_inches", 1.05))),
        )
        text_x = Inches(1.55)
    else:
        text_x = x

    _text_box(
        slide,
        text_x,
        Inches(8.22),
        Inches(6.5),
        Inches(0.3),
        theme["ppt"].get("footer_text", ""),
        size=Pt(9),
        color=colors["muted"],
        font_name=theme["fonts"].get("body"),
    )


def _resolve_content_path(value: str) -> Path:
    path = Path(value)
    if path.is_dir():
        path = path / "content.json"
    if not path.exists():
        raise SystemExit(f"content file not found: {path}")
    return path


def _default_output_path(content: dict, content_path: Path) -> Path:
    safe = "".join(
        c if c.isalnum() or c == "_" else "_"
        for c in content.get("title", "deck").lower().replace(" ", "_")
    )[:30]
    return content_path.parent / f"{safe}.pptx"


def main():
    parser = argparse.ArgumentParser(description="Generate PPTX slides from content.json")
    parser.add_argument("input", help="Path to content.json or run directory")
    parser.add_argument("--output", default=None, help="Output .pptx path")
    parser.add_argument("--theme", default=DEFAULT_THEME_NAME, help="Theme name or YAML path")
    parser.add_argument("--template", default=None, help="Optional PPTX template path")
    parser.add_argument("--logo", default=None, help="Optional logo image path")
    args = parser.parse_args()

    content_path = _resolve_content_path(args.input)
    content = json.loads(content_path.read_text(encoding="utf-8"))
    output_path = Path(args.output) if args.output else _default_output_path(content, content_path)
    try:
        generate_pptx(
            content,
            str(output_path),
            theme=args.theme,
            template_path=args.template,
            logo_path=args.logo,
        )
    except (FileNotFoundError, ValueError) as exc:
        raise SystemExit(f"ERROR: {exc}") from exc
    print(f"PPTX regenerated: {output_path}")


if __name__ == "__main__":
    main()
